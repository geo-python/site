#!/usr/bin/env python3
"""
Generate PowerPoint presentation for Ready Santa Fe
Public Safety Committee presentation - January 20, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def create_title_slide(prs):
    """Create the title slide"""
    slide_layout = prs.slide_layouts[0]  # Title Slide layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Ready Santa Fe"

    subtitle_text = """Enhanced Public Information and Warning for Santa Fe
Public Safety Committee
January 20, 2026

Kyle Morgan, Interim Director
Office of Emergency Management"""

    subtitle.text = subtitle_text

    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide


def create_bullet_slide(prs, title_text, content_items):
    """Create a slide with title and bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = title_text

    # Format title
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    # Add content
    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    for item in content_items:
        if 'level' not in item:
            item['level'] = 0

        p = text_frame.add_paragraph()
        p.text = item['text']
        p.level = item['level']
        p.font.size = Pt(18) if item['level'] == 0 else Pt(16)

        if item.get('bold', False):
            p.font.bold = True

    return slide


def create_slide_2(prs):
    """What Ready Santa Fe Does"""
    content = [
        {'text': 'Steady State Operations:', 'bold': True},
        {'text': 'Visual display of road closures and public safety notices', 'level': 1},
        {'text': 'Supports Alert Santa Fe by answering "WHERE is this happening?"', 'level': 1},
        {'text': 'Reduces repeat calls to dispatch asking for location details', 'level': 1},
        {'text': ''},
        {'text': 'Disaster Response:', 'bold': True},
        {'text': 'Instant scale-up for evacuations (READY/SET/GO orders)', 'level': 1},
        {'text': 'Shelter locations and points of distribution', 'level': 1},
        {'text': 'Multiple simultaneous incidents on one map', 'level': 1},
        {'text': ''},
        {'text': 'Preparedness Resources:', 'bold': True},
        {'text': 'Local situational awareness tools', 'level': 1},
        {'text': 'Wildland fire hazard mitigation recommendations', 'level': 1},
        {'text': 'Links to hazard information and emergency planning resources', 'level': 1},
    ]

    return create_bullet_slide(prs, "What Ready Santa Fe Does", content)


def create_slide_3(prs):
    """Current Status & Success Stories"""
    content = [
        {'text': "This is NEW - We're Still Learning:", 'bold': True},
        {'text': 'Launched training program for SFPD and SFFD leadership', 'level': 1},
        {'text': '3-week practice scenario program to build familiarity', 'level': 1},
        {'text': 'Iterating based on first responder feedback', 'level': 1},
        {'text': ''},
        {'text': 'Already Deployed:', 'bold': True},
        {'text': '2 law enforcement incidents', 'level': 1},
        {'text': '1 mass care and sheltering operation (Code Blue)', 'level': 1},
        {'text': ''},
        {'text': 'Huge Thank You to SFPD:', 'bold': True},
        {'text': 'Leadership willingness to learn new technology', 'level': 1},
        {'text': 'Adaptive approach to improving public communications', 'level': 1},
        {'text': 'Partnership in developing SOPs and best practices', 'level': 1},
    ]

    return create_bullet_slide(prs, "Current Status & Success Stories", content)


def create_slide_4(prs):
    """The Critical Context"""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    title.text = "The Critical Context"
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    body_shape = slide.placeholders[1]
    text_frame = body_shape.text_frame
    text_frame.clear()

    # Add main message
    p = text_frame.add_paragraph()
    p.text = "Ready Santa Fe Supports Improved Public Information - But It's Not Everything:"
    p.font.size = Pt(18)
    p.font.bold = True

    # Add checkmark items
    items = [
        "✓ Register for Alert Santa Fe - Direct text/email notifications",
        "✓ Enable Government Alerts - IPAWS/Wireless Emergency Alerts on your phone",
        "✓ Bookmark ready.santafenm.gov - Visual geographic information"
    ]

    for item in items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_before = Pt(6)

    # Add emphasis text
    p = text_frame.add_paragraph()
    p.text = ""
    p = text_frame.add_paragraph()
    p.text = "Redundancy is valuable - We communicate across multiple platforms to ensure you receive consistent, official information however you prefer to stay informed."
    p.font.size = Pt(16)
    p.font.italic = True
    p.space_before = Pt(12)

    return slide


def create_slide_5(prs):
    """Where We're Going"""
    content = [
        {'text': 'Near-Term Integrations:', 'bold': True},
        {'text': 'Apple Maps and Google Maps (traffic/closure data)', 'level': 1},
        {'text': 'Third-party apps like Watch Duty (wildfire evacuations)', 'level': 1},
        {'text': 'NIFC portal and other authoritative sources', 'level': 1},
        {'text': ''},
        {'text': 'Why? People interact with spatial data in many ways throughout their day. We want our information available where residents already look.', 'level': 0},
        {'text': ''},
        {'text': 'Future Growth Areas:', 'bold': True},
        {'text': 'Special events management (pre-planned closures)', 'level': 1},
        {'text': 'Utility notifications (water, power outages)', 'level': 1},
        {'text': 'Enhanced multi-agency coordination tools', 'level': 1},
        {'text': ''},
        {'text': 'The Goal: Make official emergency information accessible, visual, and scalable for any incident size.', 'level': 0, 'bold': True},
    ]

    return create_bullet_slide(prs, "Where We're Going", content)


def main():
    """Generate the Ready Santa Fe presentation"""
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create slides
    print("Creating Slide 1: Title Slide...")
    create_title_slide(prs)

    print("Creating Slide 2: What Ready Santa Fe Does...")
    create_slide_2(prs)

    print("Creating Slide 3: Current Status & Success Stories...")
    create_slide_3(prs)

    print("Creating Slide 4: The Critical Context...")
    create_slide_4(prs)

    print("Creating Slide 5: Where We're Going...")
    create_slide_5(prs)

    # Save presentation
    output_file = "Ready_Santa_Fe_Presentation.pptx"
    prs.save(output_file)
    print(f"\n✓ Presentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")
    print("\nPresentation is ready for the Public Safety Committee meeting on January 20, 2026")


if __name__ == "__main__":
    main()
